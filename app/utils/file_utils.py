import uuid
from pathlib import Path
import fitz        # PyMuPDF
import docx        # python-docx
import pptx        # python-pptx
from fastapi import HTTPException

# where to save uploads
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

# only allow these extensions
ALLOWED_EXTS = {".pdf", ".docx", ".pptx"}

def save_upload(upload_file) -> Path:
    """
    Validate extension, save file, and return Path to saved file.
    Raises HTTPException(400) if unsupported.
    """
    ext = Path(upload_file.filename).suffix.lower()
    if ext not in ALLOWED_EXTS:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

    unique_name = f"{uuid.uuid4()}{ext}"
    out_path = UPLOAD_DIR / unique_name
    with out_path.open("wb") as f:
        f.write(upload_file.file.read())
    return out_path

def extract_text(path: Path) -> str:
    """
    Extracts and returns all text from the file at `path`.
    Supports PDF, DOCX (including tables), PPTX.
    """
    ext = path.suffix.lower()

    if ext == ".pdf":
        doc = fitz.open(path)
        return "\n".join(page.get_text() for page in doc)

    if ext == ".docx":
        d = docx.Document(path)
        parts = []
        # paragraphs
        parts.extend(p.text for p in d.paragraphs)
        # tables
        for table in d.tables:
            for row in table.rows:
                for cell in row.cells:
                    # cell.text returns full text of a cell
                    parts.append(cell.text)
        return "\n".join(parts)

    if ext == ".pptx":
        prs = pptx.Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    # Should never happen if ALLOWED_EXTS is enforced
    return ""
